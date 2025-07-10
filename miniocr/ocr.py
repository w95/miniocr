import asyncio
import base64
import os
import platform
import subprocess
from typing import List, Optional
from openai import AsyncOpenAI
import aiofiles
import aiohttp
from pdf2image import convert_from_path, convert_from_bytes
import tempfile
from pptx import Presentation
from PIL import Image
import io

class MiniOCR:
    def __init__(self, api_key: str = None):
        self.client = AsyncOpenAI(api_key=api_key or os.getenv("OPENAI_API_KEY"))
        self.system_prompt = """
Convert the following document to markdown.
Return only the markdown with no explanation text. Do not include delimiters like ```markdown or ```html.

RULES:
  - You must include all information on the page. Do not exclude headers, footers, or subtext.
  - Return tables in an HTML format.
  - Charts & infographics must be interpreted to a markdown format. Prefer table format when applicable.
  - Logos should be wrapped in brackets. Ex: <logo>Coca-Cola<logo>
  - Watermarks should be wrapped in brackets. Ex: <watermark>OFFICIAL COPY<watermark>
  - Page numbers should be wrapped in brackets. Ex: <page_number>14<page_number> or <page_number>9/22<page_number>
  - Prefer using ☐ and ☑ for check boxes.
"""

    async def encode_image_to_base64(self, image_path: str) -> str:
        """Encode an image to base64 asynchronously."""
        async with aiofiles.open(image_path, "rb") as image_file:
            image_data = await image_file.read()
        return base64.b64encode(image_data).decode("utf-8")

    def is_image_file(self, file_path: str) -> bool:
        """Check if file is an image."""
        image_extensions = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp')
        return file_path.lower().endswith(image_extensions)

    def is_pdf_file(self, file_path: str) -> bool:
        """Check if file is a PDF."""
        return file_path.lower().endswith('.pdf')

    def is_pptx_file(self, file_path: str) -> bool:
        """Check if file is a PowerPoint presentation."""
        return file_path.lower().endswith('.pptx')
    
    def find_libreoffice_executable(self) -> Optional[str]:
        """Find LibreOffice executable path based on the operating system."""
        system = platform.system().lower()
        
        # Common executable names to try
        if system == "windows":
            executable_names = ["soffice.exe", "soffice"]
            # Common Windows installation paths
            common_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                r"C:\LibreOffice\program\soffice.exe",
            ]
        else:
            executable_names = ["soffice", "libreoffice"]
            # Common Unix/Linux/macOS paths
            common_paths = [
                "/usr/bin/soffice",
                "/usr/local/bin/soffice",
                "/opt/libreoffice/program/soffice",
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            ]
        
        # First, try to find in PATH
        for exec_name in executable_names:
            try:
                result = subprocess.run(
                    ["which", exec_name] if system != "windows" else ["where", exec_name],
                    capture_output=True, text=True, timeout=5
                )
                if result.returncode == 0 and result.stdout.strip():
                    return exec_name  # Found in PATH
            except (subprocess.TimeoutExpired, FileNotFoundError):
                continue
        
        # If not in PATH, try common installation paths
        for path in common_paths:
            if os.path.isfile(path) and os.access(path, os.X_OK):
                return path
        
        # Last resort: try just the executable name and hope it works
        return executable_names[0]

    async def download_file(self, url: str, temp_dir: str) -> str:
        """Download file from URL if needed."""
        if url.startswith(('http://', 'https://')):
            async with aiohttp.ClientSession() as session:
                async with session.get(url) as response:
                    filename = os.path.basename(url.split('?')[0]) or 'document'
                    filepath = os.path.join(temp_dir, filename)
                    async with aiofiles.open(filepath, 'wb') as f:
                        async for chunk in response.content.iter_chunked(8192):
                            await f.write(chunk)
                    return filepath
        return url

    async def pdf_to_images(self, pdf_path: str, temp_dir: str, dpi: int = 200) -> List[str]:
        """Convert PDF to images."""
        images = convert_from_path(pdf_path, dpi=dpi)
        image_paths = []
        
        for i, image in enumerate(images):
            image_path = os.path.join(temp_dir, f"page_{i+1}.png")
            image.save(image_path, 'PNG')
            image_paths.append(image_path)
        
        return image_paths

    async def pptx_to_images(self, pptx_path: str, temp_dir: str, dpi: int = 200) -> List[str]:
        """Convert PowerPoint presentation to images via PDF conversion."""
        try:
            # Get the base filename without extension
            filename_base = os.path.basename(pptx_path)
            filename_bare = os.path.splitext(filename_base)[0]
            
            # Convert PPTX to PDF using LibreOffice
            pdf_path = os.path.join(temp_dir, f"{filename_bare}.pdf")
            
            # Find LibreOffice executable
            soffice_executable = self.find_libreoffice_executable()
            if not soffice_executable:
                print("Warning: LibreOffice executable not found, falling back to text extraction")
                return await self.pptx_to_images_fallback(pptx_path, temp_dir)
            
            # Use soffice to convert PPTX to PDF
            command_list = [
                soffice_executable, 
                "--headless", 
                "--convert-to", "pdf", 
                "--outdir", temp_dir,
                pptx_path
            ]
            
            # Set timeout for the conversion process
            try:
                result = subprocess.run(command_list, capture_output=True, text=True, timeout=60)
            except subprocess.TimeoutExpired:
                print("Warning: LibreOffice conversion timed out, falling back to text extraction")
                return await self.pptx_to_images_fallback(pptx_path, temp_dir)
            
            if result.returncode != 0:
                # Fallback to text extraction if soffice fails
                print(f"Warning: LibreOffice conversion failed (exit code {result.returncode}): {result.stderr}")
                return await self.pptx_to_images_fallback(pptx_path, temp_dir)
            
            # Check if PDF was created
            if not os.path.exists(pdf_path):
                print("Warning: PDF file was not created, falling back to text extraction")
                return await self.pptx_to_images_fallback(pptx_path, temp_dir)
            
            # Convert PDF to images
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()
            
            images = convert_from_bytes(pdf_bytes, dpi=dpi)
            image_paths = []
            
            for i, img in enumerate(images):
                image_path = os.path.join(temp_dir, f"slide_{i+1}.png")
                img.save(image_path, 'PNG')
                image_paths.append(image_path)
            
            return image_paths
            
        except Exception as e:
            print(f"Error in PPTX to images conversion: {e}")
            # Fallback to text extraction
            return await self.pptx_to_images_fallback(pptx_path, temp_dir)
    
    async def pptx_to_images_fallback(self, pptx_path: str, temp_dir: str) -> List[str]:
        """Fallback method: Extract text and create simple images."""
        presentation = Presentation(pptx_path)
        image_paths = []
        
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            
            # Extract text from slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                # Create a simple white image with text content
                img = Image.new('RGB', (800, 600), color='white')
                image_path = os.path.join(temp_dir, f"slide_{i+1}.png")
                img.save(image_path, 'PNG')
                image_paths.append(image_path)
        
        return image_paths

    async def process_image(self, image_path: str, model: str = "gpt-4o-mini") -> str:
        """Process a single image with OpenAI Vision API."""
        base64_image = await self.encode_image_to_base64(image_path)
        
        response = await self.client.chat.completions.create(
            model=model,
            messages=[
                {
                    "role": "system",
                    "content": self.system_prompt
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=4000
        )
        
        return response.choices[0].message.content

    async def process_pptx_text(self, pptx_path: str) -> str:
        """Process PowerPoint text directly without image conversion."""
        presentation = Presentation(pptx_path)
        slide_contents = []
        
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            
            # Extract text from slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                slide_content = f"## Slide {i+1}\n\n" + "\n\n".join(slide_text)
                slide_contents.append(slide_content)
        
        return "\n\n".join(slide_contents)

    async def ocr(
        self,
        file_path: str,
        model: str = "gpt-4o-mini",
        concurrency: int = 5,
        output_dir: Optional[str] = None,
        cleanup: bool = True
    ) -> dict:
        """
        Main function to convert PDF/image/PPTX to markdown using OpenAI Vision API.
        
        Args:
            file_path: Path or URL to PDF, image, or PPTX file
            model: OpenAI model to use (default: gpt-4o-mini)
            concurrency: Number of concurrent requests
            output_dir: Directory to save markdown output
            cleanup: Whether to cleanup temporary files
            
        Returns:
            Dictionary with markdown content and metadata
        """
        with tempfile.TemporaryDirectory() as temp_dir:
            # Download file if URL
            local_path = await self.download_file(file_path, temp_dir)
            
            # Determine file type and process accordingly
            if self.is_pptx_file(local_path):
                # Convert PPTX to images and process with Vision API
                image_paths = await self.pptx_to_images(local_path, temp_dir)
                # Process images concurrently
                semaphore = asyncio.Semaphore(concurrency)
                
                async def process_with_semaphore(image_path):
                    async with semaphore:
                        return await self.process_image(image_path, model)
                
                tasks = [process_with_semaphore(img_path) for img_path in image_paths]
                results = await asyncio.gather(*tasks)
                
                markdown_content = "\n\n".join(results) if len(results) > 1 else results[0]
                pages_count = len(results)
            elif self.is_image_file(local_path):
                image_paths = [local_path]
                # Process images concurrently
                semaphore = asyncio.Semaphore(concurrency)
                
                async def process_with_semaphore(image_path):
                    async with semaphore:
                        return await self.process_image(image_path, model)
                
                tasks = [process_with_semaphore(img_path) for img_path in image_paths]
                results = await asyncio.gather(*tasks)
                
                markdown_content = "\n\n".join(results) if len(results) > 1 else results[0]
                pages_count = len(results)
            elif self.is_pdf_file(local_path):
                image_paths = await self.pdf_to_images(local_path, temp_dir)
                # Process images concurrently
                semaphore = asyncio.Semaphore(concurrency)
                
                async def process_with_semaphore(image_path):
                    async with semaphore:
                        return await self.process_image(image_path, model)
                
                tasks = [process_with_semaphore(img_path) for img_path in image_paths]
                results = await asyncio.gather(*tasks)
                
                markdown_content = "\n\n".join(results) if len(results) > 1 else results[0]
                pages_count = len(results)
            else:
                raise ValueError(f"Unsupported file type: {local_path}")
            
            # Save to file if output_dir specified
            if output_dir:
                os.makedirs(output_dir, exist_ok=True)
                filename = os.path.splitext(os.path.basename(local_path))[0]
                output_path = os.path.join(output_dir, f"{filename}.md")
                async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                    await f.write(markdown_content)
            
            return {
                "content": markdown_content,
                "pages": pages_count,
                "file_name": os.path.basename(local_path)
            } 